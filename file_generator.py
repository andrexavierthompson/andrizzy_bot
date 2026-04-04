"""
Shared file generation module for all bots.
Each function returns (bytes, filename).
"""

import io
import base64
import os
from datetime import date

import httpx


# ── Word Documents ────────────────────────────────────────────────────────────

def generate_word(title: str, sections: list, style: str = "plain",
                  bot_name: str = "Andre", filename_hint: str = "") -> tuple:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    if style == "polished":
        # Set default font
        style_obj = doc.styles["Normal"]
        font = style_obj.font
        font.name = "Calibri"
        font.size = Pt(11)

        # Title block
        title_para = doc.add_paragraph()
        title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run = title_para.add_run(title)
        run.bold = True
        run.font.size = Pt(20)
        run.font.name = "Calibri"

        sub = doc.add_paragraph()
        sub_run = sub.add_run(f"{bot_name}  ·  {date.today().strftime('%d %B %Y')}")
        sub_run.font.size = Pt(10)
        sub_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        doc.add_paragraph()  # spacer
    else:
        title_para = doc.add_paragraph()
        run = title_para.add_run(title)
        run.bold = True
        run.font.size = Pt(16)
        doc.add_paragraph()

    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")

        if heading:
            if style == "polished":
                h = doc.add_paragraph()
                h_run = h.add_run(heading)
                h_run.bold = True
                h_run.font.size = Pt(13)
                h_run.font.name = "Calibri"
                h_run.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
            else:
                h = doc.add_paragraph()
                h_run = h.add_run(heading)
                h_run.bold = True

        if body:
            p = doc.add_paragraph(body)
            if style == "polished":
                p.paragraph_format.space_after = Pt(8)
                for run in p.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(11)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    filename = _make_filename(filename_hint or title, "docx")
    return buf.read(), filename


# ── PDF ───────────────────────────────────────────────────────────────────────

def generate_pdf(title: str, sections: list, style: str = "plain",
                 bot_name: str = "Andre", filename_hint: str = "") -> tuple:
    from fpdf import FPDF

    class PDF(FPDF):
        def __init__(self, style, bot_name, title):
            super().__init__()
            self._style = style
            self._bot_name = bot_name
            self._doc_title = title

        def header(self):
            if self._style == "polished":
                self.set_fill_color(26, 26, 46)  # dark navy
                self.rect(0, 0, 210, 18, "F")
                self.set_font("Helvetica", "B", 11)
                self.set_text_color(255, 255, 255)
                self.set_xy(10, 4)
                self.cell(0, 10, self._doc_title, ln=False)
                self.set_text_color(0, 0, 0)
                self.ln(18)

        def footer(self):
            if self._style == "polished":
                self.set_y(-15)
                self.set_font("Helvetica", "I", 8)
                self.set_text_color(150, 150, 150)
                self.cell(0, 10,
                          f"{self._bot_name}  ·  {date.today().strftime('%d %B %Y')}  ·  Page {self.page_no()}",
                          align="C")

    pdf = PDF(style, bot_name, title)
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=20)

    if style == "polished":
        pdf.set_font("Helvetica", "B", 18)
        pdf.set_text_color(26, 26, 46)
        pdf.cell(0, 12, title, ln=True)
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(130, 130, 130)
        pdf.cell(0, 6, date.today().strftime("%d %B %Y"), ln=True)
        pdf.ln(6)
    else:
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, title, ln=True)
        pdf.ln(4)

    pdf.set_text_color(0, 0, 0)

    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")

        if heading:
            if style == "polished":
                pdf.set_font("Helvetica", "B", 12)
                pdf.set_text_color(26, 26, 46)
            else:
                pdf.set_font("Helvetica", "B", 12)
            pdf.cell(0, 8, heading, ln=True)
            pdf.set_text_color(0, 0, 0)

        if body:
            pdf.set_font("Helvetica", "", 10)
            pdf.multi_cell(0, 6, body)
            pdf.ln(4)

    buf = io.BytesIO()
    pdf_bytes = pdf.output()
    filename = _make_filename(filename_hint or title, "pdf")
    return bytes(pdf_bytes), filename


# ── Excel ─────────────────────────────────────────────────────────────────────

def generate_excel(title: str, headers: list, rows: list,
                   sheet_name: str = "Sheet1", style: str = "plain",
                   bot_name: str = "Andre", filename_hint: str = "") -> tuple:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name

    if style == "polished":
        # Title row
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 1))
        title_cell = ws.cell(row=1, column=1, value=title)
        title_cell.font = Font(bold=True, size=14, color="1A1A2E")
        title_cell.alignment = Alignment(horizontal="left")
        ws.row_dimensions[1].height = 22

        # Header row
        header_fill = PatternFill("solid", fgColor="1A1A2E")
        header_font = Font(bold=True, color="FFFFFF", size=10)
        for col_idx, header in enumerate(headers, start=1):
            cell = ws.cell(row=2, column=col_idx, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center")
        ws.freeze_panes = "A3"

        # Data rows with alternating fill
        light_fill = PatternFill("solid", fgColor="F2F2F2")
        for row_idx, row_data in enumerate(rows, start=3):
            fill = light_fill if row_idx % 2 == 1 else None
            for col_idx, value in enumerate(row_data, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                if fill:
                    cell.fill = fill

        # Auto-size columns
        for col_idx in range(1, len(headers) + 1):
            col_letter = get_column_letter(col_idx)
            max_len = max(
                (len(str(ws.cell(row=r, column=col_idx).value or ""))
                 for r in range(1, len(rows) + 3)),
                default=10
            )
            ws.column_dimensions[col_letter].width = min(max_len + 4, 50)
    else:
        # Plain: simple headers + data
        for col_idx, header in enumerate(headers, start=1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = Font(bold=True)
        for row_idx, row_data in enumerate(rows, start=2):
            for col_idx, value in enumerate(row_data, start=1):
                ws.cell(row=row_idx, column=col_idx, value=value)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    filename = _make_filename(filename_hint or title, "xlsx")
    return buf.read(), filename


# ── PowerPoint ────────────────────────────────────────────────────────────────

def generate_pptx(title: str, slides: list, style: str = "plain",
                  bot_name: str = "Andre", filename_hint: str = "") -> tuple:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    if style == "polished":
        bg = title_slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
        txBox = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        p.runs[0].font.size = Pt(40)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.runs[0].font.name = "Calibri"
        sub_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(0.5))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = f"{bot_name}  ·  {date.today().strftime('%d %B %Y')}"
        sub_p.runs[0].font.size = Pt(16)
        sub_p.runs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xCC)
        sub_p.runs[0].font.name = "Calibri"
    else:
        txBox = title_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.runs[0].font.size = Pt(36)
        p.runs[0].font.bold = True

    # Content slides
    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        slide_title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])

        if style == "polished":
            # Accent bar at top
            from pptx.util import Emu
            bar = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                0, 0, prs.slide_width, Inches(0.12)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
            bar.line.fill.background()

        # Slide title
        if slide_title:
            t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_title
            r = p.runs[0]
            r.font.size = Pt(28) if style == "polished" else Pt(24)
            r.font.bold = True
            if style == "polished":
                r.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
                r.font.name = "Calibri"

        # Bullets
        if bullets:
            b_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))
            tf = b_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets[:8]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {bullet}"
                r = p.runs[0]
                r.font.size = Pt(20) if style == "polished" else Pt(18)
                if style == "polished":
                    r.font.name = "Calibri"
                p.space_after = Pt(6)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = _make_filename(filename_hint or title, "pptx")
    return buf.read(), filename


# ── Local Save via Bridge ─────────────────────────────────────────────────────

async def save_to_local(filename: str, file_bytes: bytes, subfolder: str,
                        bridge_url: str, bridge_secret: str) -> bool:
    if not bridge_url:
        return False
    try:
        content_b64 = base64.b64encode(file_bytes).decode("utf-8")
        async with httpx.AsyncClient(timeout=15) as http:
            resp = await http.post(
                f"{bridge_url}/save-file",
                json={
                    "secret": bridge_secret,
                    "filename": filename,
                    "content_b64": content_b64,
                    "subfolder": subfolder
                }
            )
        return resp.status_code == 200
    except Exception:
        return False


# ── Helpers ───────────────────────────────────────────────────────────────────

def _make_filename(hint: str, ext: str) -> str:
    slug = hint.lower().strip()
    slug = "".join(c if c.isalnum() or c in "-_ " else "" for c in slug)
    slug = slug.replace(" ", "-")[:50].strip("-")
    today = date.today().strftime("%Y-%m-%d")
    return f"{slug}-{today}.{ext}"
